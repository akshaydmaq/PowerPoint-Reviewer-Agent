"""
PowerPoint Analysis Script - Extract and display all text content for review
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import json

def get_alignment_name(alignment):
    """Convert alignment enum to readable name"""
    if alignment is None:
        return "None"
    alignment_map = {
        PP_ALIGN.LEFT: "LEFT",
        PP_ALIGN.CENTER: "CENTER",
        PP_ALIGN.RIGHT: "RIGHT",
        PP_ALIGN.JUSTIFY: "JUSTIFY",
        PP_ALIGN.DISTRIBUTE: "DISTRIBUTE",
    }
    return alignment_map.get(alignment, str(alignment))

def analyze_shape(shape, shape_idx):
    """Analyze a single shape and return its properties"""
    info = {
        "index": shape_idx,
        "name": shape.name,
        "type": type(shape).__name__,
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }
    
    if shape.has_text_frame:
        info["has_text"] = True
        info["paragraphs"] = []
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            para_info = {
                "index": para_idx,
                "text": para.text,
                "alignment": get_alignment_name(para.alignment),
                "level": para.level,
                "runs": []
            }
            for run_idx, run in enumerate(para.runs):
                run_info = {
                    "text": run.text,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "size": run.font.size.pt if run.font.size else None,
                    "font_name": run.font.name,
                }
                para_info["runs"].append(run_info)
            info["paragraphs"].append(para_info)
    
    if shape.has_table:
        info["has_table"] = True
        info["table_data"] = []
        for row_idx, row in enumerate(shape.table.rows):
            row_data = []
            for cell_idx, cell in enumerate(row.cells):
                row_data.append(cell.text)
            info["table_data"].append(row_data)
    
    return info

def analyze_presentation(pptx_path):
    """Analyze the entire presentation"""
    prs = Presentation(pptx_path)
    
    analysis = {
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slide_count": len(prs.slides),
        "slides": []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "index": slide_idx + 1,
            "shapes": []
        }
        
        # Check for notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
            slide_info["notes"] = notes_text
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = analyze_shape(shape, shape_idx)
            slide_info["shapes"].append(shape_info)
        
        analysis["slides"].append(slide_info)
    
    return analysis

def print_text_content(analysis):
    """Print all text content in a readable format"""
    print("=" * 80)
    print("POWERPOINT CONTENT ANALYSIS")
    print("=" * 80)
    print(f"Total Slides: {analysis['slide_count']}")
    print(f"Slide Size: {analysis['slide_width']} x {analysis['slide_height']} EMUs")
    print()
    
    all_texts = []
    
    for slide in analysis["slides"]:
        print(f"\n{'='*80}")
        print(f"SLIDE {slide['index']}")
        print("=" * 80)
        
        if slide.get("notes"):
            print(f"  [NOTES]: {slide['notes'][:100]}...")
        
        for shape in slide["shapes"]:
            if shape.get("has_text") and shape.get("paragraphs"):
                has_content = any(p["text"].strip() for p in shape["paragraphs"])
                if has_content:
                    print(f"\n  Shape: {shape['name']} (Type: {shape['type']})")
                    print(f"  Position: left={shape['left']}, top={shape['top']}")
                    print(f"  Size: width={shape['width']}, height={shape['height']}")
                    for para in shape["paragraphs"]:
                        if para["text"].strip():
                            print(f"    [{para['alignment']}, Level {para['level']}] \"{para['text']}\"")
                            all_texts.append({
                                "slide": slide["index"],
                                "shape": shape["name"],
                                "text": para["text"]
                            })
                            if para["runs"]:
                                for run in para["runs"]:
                                    if run["text"].strip():
                                        font_info = f"Font: {run['font_name']}, Size: {run['size']}, Bold: {run['bold']}, Italic: {run['italic']}"
                                        print(f"      Run: \"{run['text']}\" | {font_info}")
            
            if shape.get("has_table"):
                print(f"\n  Table in {shape['name']}:")
                for row_idx, row in enumerate(shape["table_data"]):
                    print(f"    Row {row_idx}: {row}")
    
    return all_texts

if __name__ == "__main__":
    pptx_path = r"c:\Projects\pptxreview\Test deck_corrected.pptx"
    analysis = analyze_presentation(pptx_path)
    all_texts = print_text_content(analysis)
    
    print("\n" + "=" * 80)
    print("ALL TEXT CONTENT (for spell check)")
    print("=" * 80)
    for item in all_texts:
        print(f"Slide {item['slide']}, {item['shape']}: \"{item['text']}\"")
