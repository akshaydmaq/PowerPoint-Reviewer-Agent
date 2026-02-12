"""
PowerPoint Review Agent - An AI-powered agent for reviewing and correcting PowerPoint presentations.

This agent uses OpenAI GPT-4 to intelligently:
- Detect and fix spelling/grammar errors
- Analyze and standardize alignment
- Apply formatting corrections
- Make context-aware decisions

Requires: OPENAI_API_KEY in .env file or environment variable
"""

import os
import json
import re
from typing import Any, Callable
from dataclasses import dataclass, field
from enum import Enum
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Emu, Pt
from openai import OpenAI

# Load environment variables from .env file
load_dotenv()

# ============================================================================
# CONFIGURATION
# ============================================================================

MODEL = "gpt-4o"  # Can change to "gpt-4-turbo" or "gpt-4"
MAX_ITERATIONS = 20  # Maximum agent loop iterations
VERBOSE = True  # Print agent reasoning


# ============================================================================
# DATA STRUCTURES
# ============================================================================

@dataclass
class SlideContent:
    """Represents content from a single slide"""
    slide_number: int
    shapes: list[dict]
    notes: str = ""


@dataclass
class Correction:
    """Represents a correction to be applied"""
    slide_number: int
    shape_name: str
    original_text: str
    corrected_text: str
    correction_type: str  # "spelling", "grammar", "alignment", "formatting"
    reasoning: str


@dataclass
class AgentState:
    """Maintains the agent's state throughout execution"""
    presentation_path: str
    output_path: str
    slides_content: list[SlideContent] = field(default_factory=list)
    pending_corrections: list[Correction] = field(default_factory=list)
    applied_corrections: list[Correction] = field(default_factory=list)
    current_task: str = "analyze"
    iteration: int = 0
    is_complete: bool = False
    messages: list[dict] = field(default_factory=list)


# ============================================================================
# TOOLS - Functions the agent can call
# ============================================================================

def tool_extract_slide_content(state: AgentState) -> dict:
    """
    Extract all text content from the presentation for analysis.
    Returns structured data about each slide's content.
    """
    prs = Presentation(state.presentation_path)
    state.slides_content = []
    
    result = {"slides": [], "total_slides": len(prs.slides)}
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        slide_data = {
            "slide_number": slide_num,
            "shapes": [],
            "notes": ""
        }
        
        # Extract notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text if slide.notes_slide.notes_text_frame else ""
            slide_data["notes"] = notes_text
        
        # Extract shape content
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "text_content": []
            }
            
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        para_info = {
                            "text": para.text,
                            "level": para.level,
                            "runs": [{"text": run.text, "bold": run.font.bold, "italic": run.font.italic} 
                                    for run in para.runs if run.text.strip()]
                        }
                        shape_info["text_content"].append(para_info)
            
            if shape.has_table:
                shape_info["table"] = []
                for row in shape.table.rows:
                    row_data = [cell.text for cell in row.cells]
                    shape_info["table"].append(row_data)
            
            if shape_info["text_content"] or shape_info.get("table"):
                slide_data["shapes"].append(shape_info)
        
        result["slides"].append(slide_data)
        state.slides_content.append(SlideContent(
            slide_number=slide_num,
            shapes=slide_data["shapes"],
            notes=slide_data.get("notes", "")
        ))
    
    return result


def tool_analyze_text_for_errors(state: AgentState, slide_number: int, text: str) -> dict:
    """
    Use GPT to analyze a specific text for spelling and grammar errors.
    Returns suggested corrections with reasoning.
    """
    client = OpenAI()
    
    response = client.chat.completions.create(
        model=MODEL,
        messages=[
            {
                "role": "system",
                "content": """You are a professional proofreader. Analyze the given text for:
1. Spelling errors
2. Grammar errors
3. Punctuation issues
4. Awkward phrasing

Return JSON with this structure:
{
    "has_errors": true/false,
    "corrected_text": "the corrected text",
    "errors_found": [
        {"type": "spelling|grammar|punctuation", "original": "wrong", "correction": "right", "explanation": "why"}
    ]
}

Be conservative - only flag clear errors. Preserve technical terms and intentional stylistic choices.
Do NOT change meaning or rewrite for style."""
            },
            {
                "role": "user",
                "content": f"Analyze this text from slide {slide_number}:\n\n\"{text}\""
            }
        ],
        response_format={"type": "json_object"},
        temperature=0.1
    )
    
    result = json.loads(response.choices[0].message.content)
    return result


def tool_analyze_alignment(state: AgentState) -> dict:
    """
    Analyze alignment consistency across all slides.
    Returns recommendations for standardization.
    """
    prs = Presentation(state.presentation_path)
    
    # Collect title positions
    title_positions = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if "Title" in shape.name:
                title_positions.append({
                    "slide": slide_idx + 1,
                    "name": shape.name,
                    "left": shape.left,
                    "top": shape.top
                })
    
    # Find inconsistencies
    if not title_positions:
        return {"has_issues": False, "message": "No titles found"}
    
    # Count left positions
    left_counts = {}
    for pos in title_positions:
        left = pos["left"]
        left_counts[left] = left_counts.get(left, 0) + 1
    
    # Most common position
    most_common_left = max(left_counts.items(), key=lambda x: x[1])[0]
    
    # Find misaligned
    misaligned = [p for p in title_positions if p["left"] != most_common_left]
    
    return {
        "has_issues": len(misaligned) > 0,
        "standard_left_position": most_common_left,
        "misaligned_titles": misaligned,
        "recommendation": f"Align all titles to left position {most_common_left} EMUs"
    }


def tool_add_correction(state: AgentState, slide_number: int, shape_name: str, 
                        original_text: str, corrected_text: str, 
                        correction_type: str, reasoning: str) -> dict:
    """
    Add a correction to the pending corrections list.
    """
    correction = Correction(
        slide_number=slide_number,
        shape_name=shape_name,
        original_text=original_text,
        corrected_text=corrected_text,
        correction_type=correction_type,
        reasoning=reasoning
    )
    state.pending_corrections.append(correction)
    
    return {
        "status": "added",
        "correction": {
            "slide": slide_number,
            "shape": shape_name,
            "type": correction_type,
            "original": original_text,
            "corrected": corrected_text
        }
    }


def tool_apply_all_corrections(state: AgentState) -> dict:
    """
    Apply all pending corrections to the presentation and save it.
    """
    if not state.pending_corrections:
        return {"status": "no_corrections", "message": "No corrections to apply"}
    
    prs = Presentation(state.presentation_path)
    applied = []
    
    for correction in state.pending_corrections:
        slide_idx = correction.slide_number - 1
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            continue
        
        slide = prs.slides[slide_idx]
        
        for shape in slide.shapes:
            if shape.name == correction.shape_name:
                # Handle alignment corrections
                if correction.correction_type == "alignment":
                    try:
                        new_left = int(correction.corrected_text)
                        shape.left = new_left
                        applied.append({
                            "slide": correction.slide_number,
                            "type": "alignment",
                            "shape": correction.shape_name
                        })
                    except ValueError:
                        pass
                
                # Handle text corrections
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if correction.original_text in run.text:
                                run.text = run.text.replace(
                                    correction.original_text, 
                                    correction.corrected_text
                                )
                                applied.append({
                                    "slide": correction.slide_number,
                                    "type": correction.correction_type,
                                    "original": correction.original_text,
                                    "corrected": correction.corrected_text
                                })
    
    # Save the corrected presentation
    prs.save(state.output_path)
    
    state.applied_corrections.extend(state.pending_corrections)
    state.pending_corrections = []
    
    return {
        "status": "success",
        "corrections_applied": len(applied),
        "details": applied,
        "output_file": state.output_path
    }


def tool_mark_complete(state: AgentState) -> dict:
    """
    Mark the agent task as complete.
    """
    state.is_complete = True
    return {
        "status": "complete",
        "total_corrections": len(state.applied_corrections),
        "output_file": state.output_path
    }


# ============================================================================
# TOOL REGISTRY
# ============================================================================

TOOLS = {
    "extract_slide_content": {
        "function": tool_extract_slide_content,
        "description": "Extract all text content from the PowerPoint presentation for analysis",
        "parameters": {
            "type": "object",
            "properties": {},
            "required": []
        }
    },
    "analyze_text_for_errors": {
        "function": tool_analyze_text_for_errors,
        "description": "Use AI to analyze specific text for spelling and grammar errors",
        "parameters": {
            "type": "object",
            "properties": {
                "slide_number": {"type": "integer", "description": "The slide number"},
                "text": {"type": "string", "description": "The text to analyze"}
            },
            "required": ["slide_number", "text"]
        }
    },
    "analyze_alignment": {
        "function": tool_analyze_alignment,
        "description": "Analyze alignment consistency of titles and elements across slides",
        "parameters": {
            "type": "object",
            "properties": {},
            "required": []
        }
    },
    "add_correction": {
        "function": tool_add_correction,
        "description": "Add a correction to the pending list",
        "parameters": {
            "type": "object",
            "properties": {
                "slide_number": {"type": "integer", "description": "The slide number"},
                "shape_name": {"type": "string", "description": "The name of the shape to correct"},
                "original_text": {"type": "string", "description": "The original text or value"},
                "corrected_text": {"type": "string", "description": "The corrected text or value"},
                "correction_type": {"type": "string", "enum": ["spelling", "grammar", "alignment", "formatting"]},
                "reasoning": {"type": "string", "description": "Why this correction is needed"}
            },
            "required": ["slide_number", "shape_name", "original_text", "corrected_text", "correction_type", "reasoning"]
        }
    },
    "apply_all_corrections": {
        "function": tool_apply_all_corrections,
        "description": "Apply all pending corrections to the presentation and save",
        "parameters": {
            "type": "object",
            "properties": {},
            "required": []
        }
    },
    "mark_complete": {
        "function": tool_mark_complete,
        "description": "Mark the review task as complete",
        "parameters": {
            "type": "object",
            "properties": {},
            "required": []
        }
    }
}


def get_openai_tools():
    """Convert tool registry to OpenAI function calling format"""
    return [
        {
            "type": "function",
            "function": {
                "name": name,
                "description": tool["description"],
                "parameters": tool["parameters"]
            }
        }
        for name, tool in TOOLS.items()
    ]


# ============================================================================
# AGENT CORE
# ============================================================================

SYSTEM_PROMPT = """You are a PowerPoint Review Agent. Your job is to review and correct PowerPoint presentations.

Your workflow:
1. First, call extract_slide_content to get all the text from the presentation
2. For each piece of text, call analyze_text_for_errors to check for spelling/grammar issues
3. Call analyze_alignment to check for alignment inconsistencies
4. For each error found, call add_correction to queue the fix
5. Once all errors are identified, call apply_all_corrections to save the fixed presentation
6. Finally, call mark_complete to finish

IMPORTANT RULES:
- Only fix clear spelling and grammar errors
- Do NOT change technical terms or business jargon
- Do NOT rewrite for style - only fix actual errors
- Preserve original meaning
- Be thorough - check ALL text on ALL slides

Process each slide systematically. After analyzing all content and applying corrections, mark the task complete."""


def run_agent(presentation_path: str, output_path: str) -> dict:
    """
    Run the PowerPoint review agent.
    
    Args:
        presentation_path: Path to the input .pptx file
        output_path: Path for the corrected output file
    
    Returns:
        Summary of corrections made
    """
    client = OpenAI()
    
    # Initialize state
    state = AgentState(
        presentation_path=presentation_path,
        output_path=output_path,
        messages=[{"role": "system", "content": SYSTEM_PROMPT}]
    )
    
    # Initial user message
    state.messages.append({
        "role": "user",
        "content": f"Please review and correct the PowerPoint presentation at: {presentation_path}\nSave the corrected version to: {output_path}"
    })
    
    print("=" * 80)
    print("POWERPOINT REVIEW AGENT")
    print("=" * 80)
    print(f"Input: {presentation_path}")
    print(f"Output: {output_path}")
    print("=" * 80)
    
    # Agent loop
    while not state.is_complete and state.iteration < MAX_ITERATIONS:
        state.iteration += 1
        
        if VERBOSE:
            print(f"\n--- Iteration {state.iteration} ---")
        
        # Call the LLM
        response = client.chat.completions.create(
            model=MODEL,
            messages=state.messages,
            tools=get_openai_tools(),
            tool_choice="auto"
        )
        
        assistant_message = response.choices[0].message
        state.messages.append(assistant_message.model_dump())
        
        # Check if the model wants to call tools
        if assistant_message.tool_calls:
            for tool_call in assistant_message.tool_calls:
                tool_name = tool_call.function.name
                tool_args = json.loads(tool_call.function.arguments) if tool_call.function.arguments else {}
                
                if VERBOSE:
                    print(f"  Tool: {tool_name}")
                    if tool_args:
                        print(f"  Args: {json.dumps(tool_args, indent=2)[:200]}...")
                
                # Execute the tool
                if tool_name in TOOLS:
                    tool_fn = TOOLS[tool_name]["function"]
                    
                    # Call tool with appropriate arguments
                    if tool_name in ["extract_slide_content", "analyze_alignment", "apply_all_corrections", "mark_complete"]:
                        result = tool_fn(state)
                    elif tool_name == "analyze_text_for_errors":
                        result = tool_fn(state, tool_args["slide_number"], tool_args["text"])
                    elif tool_name == "add_correction":
                        result = tool_fn(
                            state,
                            tool_args["slide_number"],
                            tool_args["shape_name"],
                            tool_args["original_text"],
                            tool_args["corrected_text"],
                            tool_args["correction_type"],
                            tool_args["reasoning"]
                        )
                    else:
                        result = {"error": f"Unknown tool: {tool_name}"}
                    
                    if VERBOSE:
                        result_preview = json.dumps(result, indent=2)
                        if len(result_preview) > 300:
                            result_preview = result_preview[:300] + "..."
                        print(f"  Result: {result_preview}")
                else:
                    result = {"error": f"Tool not found: {tool_name}"}
                
                # Add tool result to messages
                state.messages.append({
                    "role": "tool",
                    "tool_call_id": tool_call.id,
                    "content": json.dumps(result)
                })
        
        elif assistant_message.content:
            if VERBOSE:
                print(f"  Agent: {assistant_message.content[:200]}...")
        
        # Check for completion
        if state.is_complete:
            break
    
    # Summary
    print("\n" + "=" * 80)
    print("AGENT COMPLETED")
    print("=" * 80)
    print(f"Iterations: {state.iteration}")
    print(f"Corrections Applied: {len(state.applied_corrections)}")
    
    if state.applied_corrections:
        print("\nCorrections:")
        for c in state.applied_corrections:
            print(f"  - Slide {c.slide_number}, {c.shape_name}: {c.correction_type}")
            print(f"    '{c.original_text}' -> '{c.corrected_text}'")
            print(f"    Reason: {c.reasoning}")
    
    print(f"\nOutput saved to: {state.output_path}")
    print("=" * 80)
    
    return {
        "iterations": state.iteration,
        "corrections": len(state.applied_corrections),
        "output_file": state.output_path,
        "details": [
            {
                "slide": c.slide_number,
                "type": c.correction_type,
                "original": c.original_text,
                "corrected": c.corrected_text
            }
            for c in state.applied_corrections
        ]
    }


# ============================================================================
# MAIN
# ============================================================================

def main():
    import sys
    
    # Check for API key
    if not os.environ.get("OPENAI_API_KEY"):
        print("ERROR: OPENAI_API_KEY not found")
        print("Add it to the .env file: OPENAI_API_KEY=sk-your-key-here")
        sys.exit(1)
    
    # Default paths
    input_path = r"c:\Projects\pptxreview\Test deck.pptx"
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    output_path = os.path.join(os.path.dirname(input_path), f"{base_name}_corrected.pptx")
    
    # Run the agent
    result = run_agent(input_path, output_path)
    
    return result


if __name__ == "__main__":
    main()
