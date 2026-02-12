"""
PowerPoint Correction Script - Fix spelling, grammar, and alignment issues
"""
from pptx import Presentation
from pptx.util import Emu
import re
import os

# Common spelling corrections dictionary (regex pattern -> replacement)
SPELLING_CORRECTIONS = {
    r'\biss\b': 'is',
    r'\bhavings\b': 'has',
    r'\bhaveing\b': 'having',
    r'\bissuess\b': 'issues',
    r'\btestssss\b': 'tests',
    r'\btestsss\b': 'tests',
    r'\btestss\b': 'tests',
    r'\bteh\b': 'the',
    r'\brecieve\b': 'receive',
    r'\boccured\b': 'occurred',
    r'\bseperately\b': 'separately',
    r'\bdefinate\b': 'definite',
    r'\boccassion\b': 'occasion',
    r'\buntill\b': 'until',
    r'\bwich\b': 'which',
    r'\bwih\b': 'with',
    r'\bwiht\b': 'with',
    r'\bthat\s+that\b': 'that',
    r'\bthe\s+the\b': 'the',
}

def simple_spell_check(text):
    """
    Simple spell check using regex patterns
    Returns corrected text and list of corrections made
    """
    corrections = []
    original = text
    
    # Apply all known corrections
    for pattern, replacement in SPELLING_CORRECTIONS.items():
        new_text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
        if new_text != text:
            corrections.append(f"'{pattern}' -> '{replacement}'")
            text = new_text
    
    # Fix obvious repeated character typos (like "testttt" -> "test")
    # Pattern: word ending with 3+ same characters
    repeated_pattern = r'\b(\w*?)([a-zA-Z])\2{2,}\b'
    
    def fix_repeated(match):
        prefix = match.group(1)
        char = match.group(2)
        # Common endings that should have double letters
        for ending in ['ss', 's', 'tt', 't', 'ee', 'e', 'll', 'l']:
            if ending[0] == char:
                return prefix + ending
        return prefix + char
    
    new_text = re.sub(repeated_pattern, fix_repeated, text)
    if new_text != text:
        corrections.append(f"Fixed repeated characters")
        text = new_text
    
    if text != original:
        corrections.insert(0, f"'{original}' -> '{text}'")
    
    return text, corrections

def get_title_placeholder_positions(prs):
    """Analyze title positions across all slides"""
    title_positions = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if "Title" in shape.name:
                title_positions.append({
                    "slide": slide_idx + 1,
                    "shape_name": shape.name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                })
    
    return title_positions

def find_most_common_position(positions):
    """Find the most common left position for titles"""
    if not positions:
        return None
    
    left_counts = {}
    for pos in positions:
        left = pos["left"]
        left_counts[left] = left_counts.get(left, 0) + 1
    
    # Return the most common position
    most_common = max(left_counts.items(), key=lambda x: x[1])
    return most_common[0]

def correct_presentation(input_path, output_path):
    """Main function to correct the presentation"""
    print(f"Opening: {input_path}")
    prs = Presentation(input_path)
    
    all_corrections = []
    
    # Analyze title positions
    print("\nAnalyzing title positions...")
    title_positions = get_title_placeholder_positions(prs)
    standard_left = find_most_common_position(title_positions)
    print(f"Standard title left position: {standard_left} EMUs")
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        print(f"\n--- Processing Slide {slide_num} ---")
        
        for shape in slide.shapes:
            # Fix title alignment
            if "Title" in shape.name and standard_left is not None:
                if shape.left != standard_left:
                    old_left = shape.left
                    shape.left = standard_left
                    all_corrections.append(f"Slide {slide_num}, {shape.name}: Aligned left from {old_left} to {standard_left}")
                    print(f"  Fixed alignment: {shape.name}")
            
            # Fix text content
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            original = run.text
                            corrected, corrections = simple_spell_check(run.text)
                            
                            if corrected != original:
                                run.text = corrected
                                for c in corrections:
                                    msg = f"Slide {slide_num}, {shape.name}: {c}"
                                    all_corrections.append(msg)
                                    print(f"  Text fix: {c}")
            
            # Fix table content
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            original = cell.text
                            corrected, corrections = simple_spell_check(cell.text)
                            
                            if corrected != original:
                                # Need to update the text frame
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.text, _ = simple_spell_check(run.text)
                                for c in corrections:
                                    msg = f"Slide {slide_num}, Table[{row_idx},{cell_idx}]: {c}"
                                    all_corrections.append(msg)
                                    print(f"  Table fix: {c}")
    
    # Save the corrected presentation
    print(f"\nSaving corrected presentation to: {output_path}")
    prs.save(output_path)
    
    return all_corrections

def main():
    input_path = r"c:\Projects\pptxreview\Test deck.pptx"
    
    # Generate output filename
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    output_dir = os.path.dirname(input_path)
    output_path = os.path.join(output_dir, f"{base_name}_corrected.pptx")
    
    print("=" * 80)
    print("POWERPOINT CORRECTION AGENT")
    print("=" * 80)
    
    corrections = correct_presentation(input_path, output_path)
    
    print("\n" + "=" * 80)
    print("CORRECTION SUMMARY")
    print("=" * 80)
    
    if corrections:
        print(f"Total corrections made: {len(corrections)}")
        for c in corrections:
            print(f"  - {c}")
    else:
        print("No corrections needed.")
    
    print(f"\nOutput saved to: {output_path}")
    print("=" * 80)

if __name__ == "__main__":
    main()
